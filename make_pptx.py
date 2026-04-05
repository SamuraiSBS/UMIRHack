# -*- coding: utf-8 -*-
"""Generate Flagman.pptx presentation."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import lxml.etree as etree

# ── paths ─────────────────────────────────────────────────
BASE   = Path("d:/UMIRHack") / "\u041f\u0440\u0435\u0437\u0435\u043d\u0442\u0430\u0446\u0438\u044f"
OUT    = Path("d:/UMIRHack") / "\u0424\u043b\u0430\u0433\u043c\u0430\u043d.pptx"
QR     = BASE / "qr-code.png"
LOGO   = BASE / "\u041b\u043e\u0433\u043e\u043a\u0443\u0440\u044c\u0435\u0440\u0430.jpg"

# sub-folders
BIZNES = BASE / "\u0431\u0438\u0437\u043d\u0435\u0441"
KURER  = BASE / "\u043a\u0443\u0440\u044c\u0435\u0440"
POKUP  = BASE / "\u043f\u043e\u043a\u0443\u043f\u0430\u0442\u0435\u043b\u044c"
MENU   = BASE / "\u043c\u0435\u043d\u044e"

# ── palette ───────────────────────────────────────────────
DARK   = RGBColor(0x1c, 0x1c, 0x1c)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x2e, 0x2e, 0x2e)
DGRAY  = RGBColor(0x12, 0x12, 0x12)
LGRAY  = RGBColor(0xcc, 0xcc, 0xcc)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, alpha=None):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    if alpha is not None:
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        srgb = sp._element.find(".//{%s}srgbClr" % ns)
        if srgb is not None:
            a = etree.SubElement(srgb, "{%s}alpha" % ns)
            a.set("val", str(alpha))
    return sp

def txt(slide, text, x, y, w, h,
        size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size    = Pt(size)
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    return tb

def pic(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        slide.shapes.add_picture(str(path), x, y, w)

def bg(slide, color=DARK):
    rect(slide, 0, 0, W, H, color)

def qr(slide):
    """QR badge top-left."""
    slide.shapes.add_picture(str(QR), Inches(0.17), Inches(0.17),
                              Inches(0.82), Inches(0.82))

def header(slide, label, bg_color=DGRAY, text_color=YELLOW):
    rect(slide, 0, 0, W, Inches(1.05), bg_color)
    txt(slide, label, Inches(1.25), Inches(0.15), W - Inches(1.5), Inches(0.78),
        size=36, bold=True, color=text_color)


# ══════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
# food photo with dark overlay
pic(s, MENU / "\u041c\u0430\u0440\u0433\u0430\u0440\u0438\u0442\u0430.jpg", 0, 0, Inches(5.5), H)
rect(s, 0, 0, Inches(5.5), H, DARK, alpha=52000)

txt(s, "\u0424\u041b\u0410\u0413\u041c\u0410\u041d",
    Inches(5.8), Inches(1.3), Inches(7.3), Inches(1.8),
    size=88, bold=True, color=YELLOW)
txt(s, "\u0410\u0433\u0440\u0435\u0433\u0430\u0442\u043e\u0440 \u0434\u043e\u0441\u0442\u0430\u0432\u043a\u0438 \u0435\u0434\u044b",
    Inches(5.8), Inches(3.1), Inches(7.2), Inches(0.7), size=24, color=WHITE)
rect(s, Inches(5.8), Inches(3.88), Inches(4.0), Inches(0.08), YELLOW)
txt(s, "\u0417\u0430\u043a\u0430\u0437\u044b\u0432\u0430\u0439 \u0431\u044b\u0441\u0442\u0440\u043e. \u0414\u043e\u0441\u0442\u0430\u0432\u043b\u044f\u0439 \u0432\u044b\u0433\u043e\u0434\u043d\u043e.",
    Inches(5.8), Inches(4.06), Inches(7.2), Inches(0.7),
    size=16, italic=True, color=LGRAY)
txt(s, "\u041a\u043e\u043c\u0430\u043d\u0434\u0430 SamuraiSBS  \u00b7  2026",
    Inches(5.8), Inches(5.7), Inches(7.2), Inches(0.55),
    size=16, bold=True, color=YELLOW)
qr(s)

# ══════════════════════════════════════════════════════════
# Slide 2 — Problem
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "\u041f\u0420\u041e\u0411\u041b\u0415\u041c\u0410")
qr(s)

cards = [
    ("\U0001f629", "\u0420\u0430\u0437\u0440\u043e\u0437\u043d\u0435\u043d\u043d\u043e\u0441\u0442\u044c",
     "\u0420\u0430\u0437\u043d\u044b\u0435 \u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u044f\n\u0434\u043b\u044f \u043a\u0430\u0436\u0434\u043e\u0433\u043e \u0440\u0435\u0441\u0442\u043e\u0440\u0430\u043d\u0430"),
    ("\u23f1", "\u0414\u043e\u043b\u0433\u043e \u0438 \u043d\u0435\u0443\u0434\u043e\u0431\u043d\u043e",
     "\u041d\u0435\u0442 \u0435\u0434\u0438\u043d\u043e\u0433\u043e \u043c\u0435\u0441\u0442\u0430\n\u0434\u043b\u044f \u0437\u0430\u043a\u0430\u0437\u0430 \u0435\u0434\u044b"),
    ("\U0001f4e6", "\u041d\u0435\u043f\u0440\u043e\u0437\u0440\u0430\u0447\u043d\u043e\u0441\u0442\u044c",
     "\u041a\u0443\u0440\u044c\u0435\u0440\u044b \u0438 \u0431\u0438\u0437\u043d\u0435\u0441\n\u0440\u0430\u0431\u043e\u0442\u0430\u044e\u0442 \u0432\u0441\u043b\u0435\u043f\u0443\u044e"),
]
cx = Inches(0.5)
for icon, title, desc in cards:
    rect(s, cx, Inches(1.5), Inches(3.95), Inches(5.3), GRAY)
    rect(s, cx, Inches(1.5), Inches(3.95), Inches(0.09), YELLOW)
    txt(s, icon, cx + Inches(0.1), Inches(1.75), Inches(3.75), Inches(1.1),
        size=52, align=PP_ALIGN.CENTER)
    txt(s, title, cx + Inches(0.2), Inches(2.95), Inches(3.55), Inches(0.62),
        size=22, bold=True, color=YELLOW)
    txt(s, desc, cx + Inches(0.2), Inches(3.65), Inches(3.55), Inches(1.6),
        size=16, color=WHITE)
    cx += Inches(4.35)
rect(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.07), YELLOW)

# ══════════════════════════════════════════════════════════
# Slide 3 — Solution
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "\u0420\u0415\u0428\u0415\u041d\u0418\u0415")
qr(s)

pic(s, MENU / "\u0431\u0443\u0440\u0433\u0435\u0440.jpg", Inches(7.2), Inches(1.08), Inches(6.0), Inches(6.35))
txt(s, "\u0424\u043b\u0430\u0433\u043c\u0430\u043d \u2014",
    Inches(0.5), Inches(1.3), Inches(6.4), Inches(0.95), size=44, bold=True, color=YELLOW)
txt(s, "\u0435\u0434\u0438\u043d\u0430\u044f \u043f\u043b\u0430\u0442\u0444\u043e\u0440\u043c\u0430\n\u0434\u043b\u044f \u0432\u0441\u0435\u0445 \u0443\u0447\u0430\u0441\u0442\u043d\u0438\u043a\u043e\u0432\n\u0434\u043e\u0441\u0442\u0430\u0432\u043a\u0438 \u0435\u0434\u044b",
    Inches(0.5), Inches(2.25), Inches(6.4), Inches(2.5), size=28, color=WHITE)
items = [
    "\U0001f3ea  \u0411\u0438\u0437\u043d\u0435\u0441 \u0443\u043f\u0440\u0430\u0432\u043b\u044f\u0435\u0442 \u043c\u0435\u043d\u044e \u0438 \u0437\u0430\u043a\u0430\u0437\u0430\u043c\u0438",
    "\U0001f6f5  \u041a\u0443\u0440\u044c\u0435\u0440\u044b \u0440\u0430\u0431\u043e\u0442\u0430\u044e\u0442 \u0447\u0435\u0440\u0435\u0437 \u043c\u043e\u0431\u0438\u043b\u044c\u043d\u043e\u0435 \u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u0435",
    "\U0001f6d2  \u041f\u043e\u043a\u0443\u043f\u0430\u0442\u0435\u043b\u0438 \u0437\u0430\u043a\u0430\u0437\u044b\u0432\u0430\u044e\u0442 \u0438\u0437 \u043b\u044e\u0431\u043e\u0433\u043e \u0437\u0430\u0432\u0435\u0434\u0435\u043d\u0438\u044f",
    "\U0001f4ca  \u0410\u0434\u043c\u0438\u043d \u043a\u043e\u043d\u0442\u0440\u043e\u043b\u0438\u0440\u0443\u0435\u0442 \u0432\u0441\u044e \u043f\u043b\u0430\u0442\u0444\u043e\u0440\u043c\u0443",
]
by = Inches(4.75)
for item in items:
    txt(s, item, Inches(0.5), by, Inches(6.4), Inches(0.52),
        size=16, color=LGRAY)
    by += Inches(0.55)

# ══════════════════════════════════════════════════════════
# Slide 4 — Tech stack
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "\u0421\u0422\u0415\u041a \u0422\u0415\u0425\u041d\u041e\u041b\u041e\u0413\u0418\u0419")
qr(s)

tech = [
    ("BACKEND",  ["Node.js + Express", "PostgreSQL", "Prisma ORM", "JWT Auth"]),
    ("FRONTEND", ["React 18", "React Router v6", "Axios", "Leaflet Maps", "Vite"]),
    ("MOBILE",   ["Android (APK)", "\u041c\u043e\u0431\u0438\u043b\u044c\u043d\u044b\u0439 \u043a\u043b\u0438\u0435\u043d\u0442", "\u0434\u043b\u044f \u043a\u0443\u0440\u044c\u0435\u0440\u043e\u0432"]),
    ("INFRA",    ["Docker Compose", "REST API", "Nginx proxy"]),
]
cx = Inches(0.35)
for title, items in tech:
    rect(s, cx, Inches(1.38), Inches(3.1), Inches(5.8), GRAY)
    rect(s, cx, Inches(1.38), Inches(3.1), Inches(0.62), YELLOW)
    txt(s, title, cx + Inches(0.12), Inches(1.4), Inches(2.88), Inches(0.6),
        size=14, bold=True, color=DARK)
    ty = Inches(2.18)
    for item in items:
        txt(s, "\u25b8  " + item, cx + Inches(0.15), ty, Inches(2.82), Inches(0.52),
            size=15, color=WHITE)
        ty += Inches(0.58)
    cx += Inches(3.3)

# ══════════════════════════════════════════════════════════
# Slide 5 — 4 roles
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "4 \u0420\u041e\u041b\u0418 \u2014 1 \u041f\u041b\u0410\u0422\u0424\u041e\u0420\u041c\u0410")
qr(s)

roles = [
    ("\U0001f6d2", "\u041f\u041e\u041a\u0423\u041f\u0410\u0422\u0415\u041b\u042c",
     "\u0412\u044b\u0431\u0438\u0440\u0430\u0435\u0442 \u0440\u0435\u0441\u0442\u043e\u0440\u0430\u043d\u044b\n\u0424\u043e\u0440\u043c\u0438\u0440\u0443\u0435\u0442 \u043a\u043e\u0440\u0437\u0438\u043d\u0443\n\u041e\u0442\u0441\u043b\u0435\u0436\u0438\u0432\u0430\u0435\u0442 \u0437\u0430\u043a\u0430\u0437"),
    ("\U0001f3ea", "\u0411\u0418\u0417\u041d\u0415\u0421",
     "\u0423\u043f\u0440\u0430\u0432\u043b\u044f\u0435\u0442 \u043c\u0435\u043d\u044e\n\u0422\u043e\u0447\u043a\u0438 \u043f\u0440\u043e\u0434\u0430\u0436\n\u0410\u043d\u0430\u043b\u0438\u0442\u0438\u043a\u0430"),
    ("\U0001f6f5", "\u041a\u0423\u0420\u042c\u0415\u0420",
     "\u041c\u043e\u0431\u0438\u043b\u044c\u043d\u043e\u0435 \u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u0435\n\u0410\u043a\u0442\u0438\u0432\u043d\u044b\u0435 \u0441\u043c\u0435\u043d\u044b\n\u041c\u0430\u0440\u0448\u0440\u0443\u0442 \u043d\u0430 \u043a\u0430\u0440\u0442\u0435"),
    ("\U0001f4ca", "\u0410\u0414\u041c\u0418\u041d",
     "\u041c\u043e\u0434\u0435\u0440\u0430\u0446\u0438\u044f\n\u0421\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430\n\u0411\u043b\u043e\u043a\u0438\u0440\u043e\u0432\u043a\u0438"),
]
cx = Inches(0.4)
for icon, role, desc in roles:
    rect(s, cx, Inches(1.45), Inches(3.0), Inches(5.7), GRAY)
    txt(s, icon, cx + Inches(0.1), Inches(1.65), Inches(2.8), Inches(1.1),
        size=52, align=PP_ALIGN.CENTER)
    txt(s, role, cx + Inches(0.1), Inches(2.8), Inches(2.8), Inches(0.58),
        size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(s, desc, cx + Inches(0.1), Inches(3.45), Inches(2.8), Inches(2.2),
        size=15, color=WHITE, align=PP_ALIGN.CENTER)
    cx += Inches(3.25)

# ══════════════════════════════════════════════════════════
# Slide 6 — Customer screenshots
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, Inches(1.05), YELLOW)
txt(s, "\U0001f6d2  \u041f\u041e\u041a\u0423\u041f\u0410\u0422\u0415\u041b\u042c",
    Inches(1.25), Inches(0.12), W - Inches(1.5), Inches(0.82),
    size=34, bold=True, color=DARK)
qr(s)
cx = Inches(0.2)
for fname in [
    "\u0433\u043b\u0430\u0432\u043d\u043e\u0435\u043c\u0435\u043d\u044e.jpg",
    "\u0440\u0435\u0441\u0442\u043e\u0440\u0430\u043d.jpg",
    "\u043a\u043e\u0440\u0437\u0438\u043d\u0430.jpg",
    "\u043c\u043e\u0438\u0437\u0430\u043a\u0430\u0437\u044b.jpg",
]:
    pic(s, POKUP / fname, cx, Inches(1.17), Inches(3.2), Inches(6.18))
    cx += Inches(3.3)

# ══════════════════════════════════════════════════════════
# Slide 7 — Business screenshots
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, Inches(1.05), YELLOW)
txt(s, "\U0001f3ea  \u0411\u0418\u0417\u041d\u0415\u0421",
    Inches(1.25), Inches(0.12), W - Inches(1.5), Inches(0.82),
    size=34, bold=True, color=DARK)
qr(s)
cx = Inches(0.2)
for fname in [
    "\u0433\u043b\u0430\u0432\u043d\u0430\u044f.jpg",
    "\u043c\u0435\u043d\u044e.jpg",
    "\u0430\u043d\u0430\u043b\u0438\u0442\u0438\u043a\u0430.jpg",
    "\u043d\u0430\u0441\u0442\u0440\u043e\u0439\u043a\u0430\u0437\u0430\u0432\u0435\u0434\u0435\u043d\u0438\u044f.jpg",
]:
    pic(s, BIZNES / fname, cx, Inches(1.17), Inches(3.2), Inches(6.18))
    cx += Inches(3.3)

# ══════════════════════════════════════════════════════════
# Slide 8 — Courier screenshots
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, Inches(1.05), YELLOW)
txt(s, "\U0001f6f5  \u041a\u0423\u0420\u042c\u0415\u0420 \u2014 \u041c\u041e\u0411\u0418\u041b\u042c\u041d\u041e\u0415 \u041f\u0420\u0418\u041b\u041e\u0416\u0415\u041d\u0418\u0415",
    Inches(1.25), Inches(0.12), W - Inches(1.5), Inches(0.82),
    size=30, bold=True, color=DARK)
qr(s)
pic(s, LOGO, Inches(0.15), Inches(1.25), Inches(2.3), Inches(2.3))
cx = Inches(2.65)
for fname in [
    "\u041f\u0430\u043d\u0435\u043b\u044c.jpg",
    "\u043c\u043e\u0438\u0437\u0430\u043a\u0430\u0437\u044b.jpg",
    "\u043a\u0430\u0440\u0442\u0430.jpg",
    "\u0432\u043f\u0443\u0442\u0438\u043a\u043a\u043b\u0438\u0435\u043d\u0442\u0443.jpg",
    "\u0438\u0441\u0442\u043e\u0440\u0438\u044f.jpg",
]:
    pic(s, KURER / fname, cx, Inches(1.17), Inches(2.1), Inches(6.18))
    cx += Inches(2.18)

# ══════════════════════════════════════════════════════════
# Slide 9 — Key features
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "\u041a\u041b\u042e\u0427\u0415\u0412\u042b\u0415 \u0424\u0418\u0427\u0418")
qr(s)

feats = [
    ("\U0001f5fa", "\u041a\u0430\u0440\u0442\u0430 \u0434\u043e\u0441\u0442\u0430\u0432\u043a\u0438",
     "Leaflet + \u0433\u0435\u043e\u043b\u043e\u043a\u0430\u0446\u0438\u044f\n\u0438 \u043c\u0430\u0440\u0448\u0440\u0443\u0442\u044b \u0432 \u0440\u0435\u0430\u043b\u044c\u043d\u043e\u043c \u0432\u0440\u0435\u043c\u0435\u043d\u0438"),
    ("\u26a1", "\u0420\u0435\u0430\u043b\u044c\u043d\u043e\u0435 \u0432\u0440\u0435\u043c\u044f",
     "\u0421\u0442\u0430\u0442\u0443\u0441 \u0437\u0430\u043a\u0430\u0437\u0430\n\u043e\u0431\u043d\u043e\u0432\u043b\u044f\u0435\u0442\u0441\u044f \u043c\u0433\u043d\u043e\u0432\u0435\u043d\u043d\u043e"),
    ("\U0001f510", "JWT \u0430\u0432\u0442\u043e\u0440\u0438\u0437\u0430\u0446\u0438\u044f",
     "\u0420\u043e\u043b\u0435\u0432\u0430\u044f \u043c\u043e\u0434\u0435\u043b\u044c\n4 \u0442\u0438\u043f\u0430 \u043f\u043e\u043b\u044c\u0437\u043e\u0432\u0430\u0442\u0435\u043b\u0435\u0439"),
    ("\U0001f4ca", "\u0410\u043d\u0430\u043b\u0438\u0442\u0438\u043a\u0430",
     "\u0421\u0442\u0430\u0442\u0438\u0441\u0442\u0438\u043a\u0430 \u043f\u0440\u043e\u0434\u0430\u0436\n\u0434\u043b\u044f \u0431\u0438\u0437\u043d\u0435\u0441\u0430 \u0438 \u0430\u0434\u043c\u0438\u043d\u0430"),
    ("\U0001f4f1", "\u041c\u043e\u0431\u0438\u043b\u044c\u043d\u044b\u0439 \u043a\u0443\u0440\u044c\u0435\u0440",
     "APK-\u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u0435\n\u0434\u043b\u044f Android"),
    ("\U0001f433", "Docker",
     "\u041e\u0434\u0438\u043d docker-compose\n\u0434\u043b\u044f \u0432\u0441\u0435\u0433\u043e \u0441\u0442\u0435\u043a\u0430"),
]
for i, (icon, title, desc) in enumerate(feats):
    fy = Inches(1.35) if i < 3 else Inches(4.22)
    fx = Inches(0.35) + (i % 3) * Inches(4.3)
    rect(s, fx, fy, Inches(4.12), Inches(2.62), GRAY)
    rect(s, fx, fy, Inches(0.62), Inches(2.62), YELLOW)
    txt(s, icon,  fx + Inches(0.75), fy + Inches(0.1), Inches(3.2), Inches(0.75), size=34)
    txt(s, title, fx + Inches(0.75), fy + Inches(0.88), Inches(3.2), Inches(0.52),
        size=17, bold=True, color=YELLOW)
    txt(s, desc,  fx + Inches(0.75), fy + Inches(1.42), Inches(3.2), Inches(1.1),
        size=14, color=WHITE)

# ══════════════════════════════════════════════════════════
# Slide 10 — Food showcase
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "\u0427\u0422\u041e \u0417\u0410\u041a\u0410\u0417\u042b\u0412\u0410\u042e\u0422")
qr(s)

food = [
    ("\u043f\u0435\u043f\u0435\u0440\u043e\u043d\u0438.jpg",     "\u041f\u0435\u043f\u043f\u0435\u0440\u043e\u043d\u0438"),
    ("\u0441\u043b\u043e\u0441\u043e\u0441\u0435\u043c.jpg",    "\u0421\u0443\u0448\u0438 \u0441 \u043b\u043e\u0441\u043e\u0441\u0435\u043c"),
    ("\u0442\u0430\u043a\u043e.jpg",                             "\u0422\u0430\u043a\u043e"),
    ("\u0447\u0438\u0437\u043a\u0435\u0439\u043a.jpg",          "\u0427\u0438\u0437\u043a\u0435\u0439\u043a"),
    ("\u043a\u0430\u0440\u0442\u043e\u0448\u043a\u0430\u0444\u0440\u0438.jpg", "\u041a\u0430\u0440\u0442\u043e\u0448\u043a\u0430 \u0444\u0440\u0438"),
    ("\u041c\u0430\u0440\u0433\u0430\u0440\u0438\u0442\u0430.jpg", "\u041c\u0430\u0440\u0433\u0430\u0440\u0438\u0442\u0430"),
]
for i, (fname, label) in enumerate(food):
    fy = Inches(1.2) if i < 3 else Inches(4.07)
    fx = Inches(0.3) + (i % 3) * Inches(4.35)
    pic(s, MENU / fname, fx, fy, Inches(4.12), Inches(2.67))
    rect(s, fx, fy + Inches(2.12), Inches(4.12), Inches(0.55), DARK)
    txt(s, label, fx + Inches(0.12), fy + Inches(2.14), Inches(3.9), Inches(0.52),
        size=15, bold=True, color=YELLOW)

# ══════════════════════════════════════════════════════════
# Slide 11 — QR finale
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, Inches(1.08), YELLOW)
txt(s, "\u041f\u041e\u041f\u0420\u041e\u0411\u0423\u0419\u0422\u0415 \u0421\u0410\u041c\u0418",
    Inches(0.5), Inches(0.13), W - Inches(1.0), Inches(0.83),
    size=42, bold=True, color=DARK, align=PP_ALIGN.CENTER)
qr(s)

for qx, label, sub in [
    (Inches(1.5),  "\U0001f310  \u0412\u0435\u0431-\u0441\u0430\u0439\u0442",
                   "\u041e\u0442\u043a\u0440\u044b\u0442\u044c \u043f\u043b\u0430\u0442\u0444\u043e\u0440\u043c\u0443"),
    (Inches(7.3),  "\U0001f4f1  APK \u041a\u0443\u0440\u044c\u0435\u0440",
                   "\u0421\u043a\u0430\u0447\u0430\u0442\u044c \u043f\u0440\u0438\u043b\u043e\u0436\u0435\u043d\u0438\u0435"),
]:
    rect(s, qx, Inches(1.25), Inches(4.5), Inches(5.8), GRAY)
    slide_qr = s.shapes.add_picture(
        str(QR), qx + Inches(0.5), Inches(1.62), Inches(3.5), Inches(3.5))
    txt(s, label, qx + Inches(0.1), Inches(5.27), Inches(4.3), Inches(0.68),
        size=22, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(s, sub,   qx + Inches(0.1), Inches(5.95), Inches(4.3), Inches(0.6),
        size=15, color=WHITE, align=PP_ALIGN.CENTER)

prs.save(str(OUT))
print("Saved:", OUT)
